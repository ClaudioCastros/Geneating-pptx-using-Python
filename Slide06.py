from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def slide06(prs, somente_titulo):
    ############## SLIDE 6

    slide = prs.slides.add_slide(somente_titulo)
    shapes = slide.shapes

    shapes.title.text = 'Adding an AutoShape'

    left = Inches(0.93)  # 0.93" centers this overall set of shapes
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)

    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = 'Step 1'

    left = left + width - Inches(0.4)
    width = Inches(2.0)  # chevrons need more width for visual balance

    for n in range(2, 6):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = 'Step %d' % n
        left = left + width - Inches(0.4)