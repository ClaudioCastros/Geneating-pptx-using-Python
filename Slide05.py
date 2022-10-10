from pptx.util import Inches, Pt
def slide05(prs, em_branco):
    ############## SLIDE 5

    img_path = 'Geneating pptx using Python\\images\\download.png'

    slide = prs.slides.add_slide(em_branco)

    left = top = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top)

    left = Inches(5)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)