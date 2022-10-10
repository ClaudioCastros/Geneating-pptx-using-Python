from pptx.util import Inches, Pt
def slide04(prs, em_branco):
    ############## SLIDE 4

    slide = prs.slides.add_slide(em_branco)

    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "This is text inside a textbox"

    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    p.font.size = Pt(40)